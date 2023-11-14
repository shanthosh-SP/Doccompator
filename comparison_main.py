import os
import subprocess
import pandas as pd
from tika import parser
from docx import Document
from tabulate import tabulate
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import pytesseract
from inscriptis import get_text
from pptx import Presentation
from bs4 import BeautifulSoup
import re
import nltk
import torch
from sklearn.metrics.pairwise import cosine_similarity
from transformers import BertTokenizer, BertModel
import logging

def ensure_directory_exists(file_path):
    directory = os.path.join(os.getcwd(), file_path)
    if not os.path.exists(directory):
        os.makedirs(directory)


logging.basicConfig(filename="comparison.log", level=logging.INFO)

def extract_text_from_pdf(input_pdf_file, output_text_file):
    # Use pdftotext with layout preservation
    subprocess.run(['pdftotext', '-layout', input_pdf_file, output_text_file])

    # Use Tika to extract page numbers from PDF metadata
    parsed_pdf = parser.from_file(input_pdf_file)
    num_pages = int(parsed_pdf['metadata'].get('xmpTPg:NPages', 0))

    # Function to extract text with layout preservation
    def extract_text_with_layout(file_path):
        with open(file_path, 'r', encoding='utf-8') as text_file:
            text = text_file.read()
        return text

    # Read the extracted text with layout preservation
    extracted_text = extract_text_with_layout(output_text_file)

    # Split the text into pages based on page breaks
    pages = extracted_text.split('\x0c')

    # Create a list to store the page numbers along with their corresponding text
    pages_with_numbers = []

    for page_num, page_text in enumerate(pages, 1):
        if page_text.strip():  # Check if the page is not empty or mostly empty
            pages_with_numbers.append(f"Page {page_num}\n{page_text}")

    # Combine all pages
    combined_text = '\n'.join(pages_with_numbers)

    # Create a new text file that combines page numbers and text
    with open(output_text_file, 'w', encoding='utf-8') as combined_text_file:
        combined_text_file.write(combined_text)

    print(f"Combined text with page numbers and layout for PDF saved to {output_text_file}")

def extract_text_from_excel(input_excel_file):
    # Create an empty DataFrame to store the data from all sheets
    all_sheets_data = pd.DataFrame()

    # Read all sheets from the Excel file into a dictionary of DataFrames
    xls = pd.ExcelFile(input_excel_file)
    sheets_data = {sheet_name: xls.parse(sheet_name) for sheet_name in xls.sheet_names}

    def save_text_to_file(sheets_data, output_file):
        # Create or open the text file for writing with 'utf-8' encoding
        with open(output_file, 'w', encoding='utf-8') as text_file:
            # Iterate through all sheets and write data to the text file
            for sheet_name, sheet_data in sheets_data.items():
                text_file.write(f"Sheet Name: {sheet_name}\n")
                sheet_data = sheet_data.fillna('')
                text_file.write(sheet_data.to_string(index=False, header=True))
                text_file.write("\n\n")

    input_excel_file = os.path.splitext(input_excel_file)[0].strip()
    output_excel_file = f"{input_excel_file}_extracted.txt"

    if sheets_data:
        save_text_to_file(sheets_data, output_excel_file)
        print(f"Extracted text for Excel saved to {output_excel_file}")
    return output_excel_file

def extract_text_from_docx(input_docs_file):
    table_number = 1

    last_content_was_table = False
    # Create a Document object to open and read the Word document
    doc = Document(input_docs_file)
    # Initialize a list to store the extracted content (text and tables) in the correct order
    document_content = []
    # Initialize a flag to indicate whether the last content extracted was a table
    last_content_was_table = False
    # Iterate through the paragraphs and tables in the document
    for paragraph in doc.paragraphs:
        # Extract and append text from the paragraph
        paragraph_text = paragraph.text.strip()
        if paragraph_text:
            if last_content_was_table:
                # If the last content was a table, add a newline before the paragraph
                document_content.append("\n" + paragraph_text)
            else:
                document_content.append(paragraph_text)
            last_content_was_table = False

    for table in doc.tables:
        table_data = {
            "table_number": table_number,
            "table_content": []
        }

        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            table_data["table_content"].append(row_data)

        # Append the table data to the document content
        document_content.append(table_data)

        table_number += 1
        last_content_was_table = True

    # Define the path for the output text file
    input_docs_file = os.path.splitext(input_docs_file)[0].strip()
    output_docs_file = f'{input_docs_file}_extracted.txt'

    # Write the extracted content to the text file while preserving the order
    with open(output_docs_file, "w", encoding="utf-8") as txt_file:
        for content in document_content:
            if isinstance(content, str):
                txt_file.write(content + "\n")  # Write text
            elif isinstance(content, dict) and "table_number" in content and "table_content" in content:
                # This is a table
                table_number = content["table_number"]
                table_data = content["table_content"]

                # Format and write the table using tabulate
                txt_file.write(f"Table {table_number}:\n")
                table_str = tabulate(table_data, tablefmt="grid")
                txt_file.write(table_str + "\n\n")

    print(f"Extracted content for Word document saved to {output_docs_file}")
    return output_docs_file

def extract_text_from_ppt(pptx_file):
    def extract_table(table):
        data = []
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                cell_text = ""
                for paragraph in cell.text_frame.paragraphs:
                    cell_text += " ".join(run.text for run in paragraph.runs)
                row_text.append(cell_text)
            data.append("\n".join(row_text))
        return "\n".join(data)

    def extract_text_from_image(image_path):
        try:
            # Open the image using PIL
            image = Image.open(image_path)

            # Use Tesseract to perform OCR on the image and extract text
            extracted_text = pytesseract.image_to_string(image)

            return extracted_text
        except Exception as e:
            print(f"Error while extracting text from image: {e}")
            return ""

    extracted_text = ""
    prs = Presentation(pptx_file)
    input_pptx_file = os.path.splitext(pptx_file)[0].strip()
    output_folder = f'{input_pptx_file}_images'

    for slide_number, slide in enumerate(prs.slides):
        extracted_text += f"Slide {slide_number + 1}:\n"

        for shape in slide.shapes:
            if shape.has_table:
                table_text = extract_table(shape.table)
                extracted_text += table_text + '\n\n'
            elif hasattr(shape, "text"):
                extracted_text += shape.text + '\n\n'
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Get the image part
                image_part = shape.image
                if image_part:
                    # Get image "file" contents
                    image_bytes = image_part.blob
                    # Make up a name for the file, e.g., 'image.jpg'
                    image_filename = os.path.join(output_folder, f'image_slide_{slide_number + 1}.jpg')
                    with open(image_filename, 'wb') as f:
                        f.write(image_bytes)
                    # Extract text from the image
                    extracted_image_text = extract_text_from_image(image_filename)
                    extracted_text += extracted_image_text + '\n\n'

    input_pptx = os.path.splitext(pptx_file)[0].strip()
    output_text_file = f'{input_pptx}_ppt_text.txt'
    with open(output_text_file, 'w', encoding='utf-8') as text_file:
        text_file.write(extracted_text)            

    return extracted_text


def extract_text_from_html(input_HTML_path_file, html_text):
    # Read the HTML content from the file
    with open(input_HTML_path_file, 'r', encoding='ISO-8859-1') as file:
        html_content = file.read()

    # Extract text content from the HTML
    text_content = get_text(html_content)

    # Save the extracted text to a text file
    with open(html_text, 'w', encoding='utf-8') as output_file:
        output_file.write(text_content)
    return html_text



# Initialize the NLTK tokenizer
nltk.download("punkt")

def compare_pdf_with_html(pdf_text, html_text,input_pdf_file):
    # Tokenize the text
    
    pdf_words = set(nltk.word_tokenize(pdf_text))
    html_words = set(nltk.word_tokenize(html_text))

    # Finding the difference
    difference_words = pdf_words - html_words

    # Finding the line, position, page
    word_positions = {}
    pdf_lines = pdf_text.splitlines()

    page_number = 0
    line_number = 0

    for line in pdf_lines:
        if line.startswith("Page"):
            match = re.match(r'Page (\d+)', line)
            if match:
                page_number = int(match.group(1))
            line_number = 0
        else:
            line_number += 1

        line_words = list(re.findall(r'\b\w+\b', line))
        for word in line_words:
            if word in difference_words:
                if word not in word_positions:
                    word_positions[word] = []
                word_positions[word].append({
                    'Page': page_number,
                    'Line': line_number,
                    'Position': line_words.index(word),
                    'LineContent': line
                })

    output = f"{input_pdf_file}_extracted.txt"
    with open(output, 'w', encoding='utf-8') as result_file:
        for word, positions in word_positions.items():
            for data in positions:
                result_file.write(f"Word: {word}\n")
                result_file.write(f"Page: {data['Page']}\n")
                result_file.write(f"Line: {data['Line']}\n")
                result_file.write(f"Position: {data['Position']}\n")
                result_file.write(f"Line Content: {data['LineContent']}\n\n")

    logging.info(f"Words in PDF but not in HTML, along with positions, saved to {output}")

    # Comparison using Jaccard Similarity
    jaccard_similarity = len(pdf_words.intersection(html_words)) / len(pdf_words.union(html_words))
    percentage_difference = (1 - jaccard_similarity) * 100
    logging.info(f"Jaccard Similarity Score: {jaccard_similarity:.2f}")

    # Comparison using BERT Cosine Similarity
    model_name = "bert-base-uncased"
    tokenizer = BertTokenizer.from_pretrained(model_name)
    model = BertModel.from_pretrained(model_name)

    tokens1 = tokenizer(pdf_text, return_tensors='pt', padding=True, truncation=True)
    tokens2 = tokenizer(html_text, return_tensors='pt', padding=True, truncation=True)

    with torch.no_grad():
        embeddings1 = model(**tokens1).last_hidden_state.mean(dim=1)
        embeddings2 = model(**tokens2).last_hidden_state.mean(dim=1)

    similarity = cosine_similarity(embeddings1, embeddings2)
    logging.info(f"Bert Cosine Similarity: {similarity[0][0]:.2f}")

def compare_excel_with_html(excel_text, html_text,input_excel_file):
    # Tokenize words
    excel_words = set(nltk.word_tokenize(excel_text))
    html_words = set(nltk.word_tokenize(html_text))

    # Calculate Jaccard similarity
    jaccard_similarity = len(excel_words.intersection(html_words)) / len(excel_words.union(html_words))
    percentage_difference = (1 - jaccard_similarity) * 100
    logging.info(f"Jaccard Similarity Score: {jaccard_similarity:.2f}")

    # Find difference words
    difference_words = excel_words - html_words

    # Find the line, position, page
    word_positions = {}
    excel_lines = excel_text.splitlines()

    sheet_number = 0
    line_number = 0

    for line in excel_lines:
        if line.startswith("sheet name"):
            match = re.match(r'sheet name: sheet(\d+)', line)
            if match:
                sheet_number = int(match.group(1))
            line_number = 0
        else:
            line_number += 1

        line_words = list(re.findall(r'\b\w+\b', line))
        for word in line_words:
            if word in difference_words:
                if word not in word_positions:
                    word_positions[word] = []
                word_positions[word].append({
                    'Sheet': sheet_number,
                    'Line': line_number,
                    'Position': line_words.index(word),
                    'LineContent': line
                })

    output = f"{input_excel_file}_extracted.txt"
    with open(output, 'w', encoding='utf-8') as result_file:
        for word, positions in word_positions.items():
            for data in positions:
                result_file.write(f"Word: {word}\n")
                result_file.write(f"Sheet: {data['Sheet']}\n")
                result_file.write(f"Line: {data['Line']}\n")
                result_file.write(f"Position: {data['Position']}\n")
                result_file.write(f"Line Content: {data['LineContent']}\n\n")

    logging.info(f"Words in excel but not in HTML, along with positions, saved to {output}")

    # Calculate cosine similarity using BERT embeddings
    model_name = "bert-base-uncased"
    tokenizer = BertTokenizer.from_pretrained(model_name)
    model = BertModel.from_pretrained(model_name)

    tokens1 = tokenizer(excel_text, return_tensors='pt', padding=True, truncation=True)
    tokens2 = tokenizer(html_text, return_tensors='pt', padding=True, truncation=True)

    with torch.no_grad():
        embeddings1 = model(**tokens1).last_hidden_state.mean(dim=1)
        embeddings2 = model(**tokens2).last_hidden_state.mean(dim=1)

    similarity = cosine_similarity(embeddings1, embeddings2)
    logging.info(f"Bert Cosine Similarity: {similarity[0][0]:.2f}")


def compare_docx_with_html(docx_text, html_text,input_docs_file):
    # Tokenize the text
    docs_words = set(nltk.word_tokenize(docx_text))
    html_words = set(nltk.word_tokenize(html_text))

    # Finding the difference
    difference_words = docs_words - html_words

    # Finding the line, position, page
    word_positions = {}
    docs_lines = docx_text.splitlines()

    page_number = 0
    line_number = 0

    for line in docs_lines:
        if line.startswith("Page"):
            match = re.match(r'Page (\d+)', line)
            if match:
                page_number = int(match.group(1))
            line_number = 0
        else:
            line_number += 1

        line_words = list(re.findall(r'\b\w+\b', line))
        for word in line_words:
            if word in difference_words:
                if word not in word_positions:
                    word_positions[word] = []
                word_positions[word].append({
                    'Page': page_number,
                    'Line': line_number,
                    'Position': line_words.index(word),
                    'LineContent': line
                })

    output = f'{input_docs_file}_docs.txt'
    with open(output, 'w', encoding='utf-8') as result_file:
        for word, positions in word_positions.items():
            for data in positions:
                result_file.write(f"Word: {word}\n")
                result_file.write(f"Page: {data['Page']}\n")
                result_file.write(f"Line: {data['Line']}\n")
                result_file.write(f"Position: {data['Position']}\n")
                result_file.write(f"Line Content: {data['LineContent']}\n\n")

    logging.info(f"Words in Docx but not in HTML, along with positions, saved to {output}")

    # Comparison using Jaccard Similarity
    jaccard_similarity = len(docs_words.intersection(html_words)) / len(docs_words.union(html_words))
    percentage_difference = (1 - jaccard_similarity) * 100
    logging.info(f"Jaccard Similarity Score: {jaccard_similarity:.2f}")

    # Comparison using BERT Cosine Similarity
    model_name = "bert-base-uncased"
    tokenizer = BertTokenizer.from_pretrained(model_name)
    model = BertModel.from_pretrained(model_name)

    tokens1 = tokenizer(docx_text, return_tensors='pt', padding=True, truncation=True)
    tokens2 = tokenizer(html_text, return_tensors='pt', padding=True, truncation=True)

    with torch.no_grad():
        embeddings1 = model(**tokens1).last_hidden_state.mean(dim=1)
        embeddings2 = model(**tokens2).last_hidden_state.mean(dim=1)

    similarity = cosine_similarity(embeddings1, embeddings2)
    logging.info(f"Bert Cosine Similarity: {similarity[0][0]:.2f}")

def compare_ppt_with_html(ppt_text, html_text,input_pptx_file):
    # Tokenize the text
    ppt_words = set(nltk.word_tokenize(ppt_text))
    html_words = set(nltk.word_tokenize(html_text))

    # Finding the difference
    difference_words = ppt_words - html_words

    # Finding the slide, line, position
    word_positions = {}
    ppt_lines = ppt_text.splitlines()

    slide_number = 0
    line_number = 0

    for line in ppt_lines:
        if line.startswith("Slide"):
            match = re.match(r'Slide (\d+):', line)
            if match:
                slide_number = int(match.group(1))
            line_number = 0
        else:
            line_number += 1

        line_words = list(re.findall(r'\b\w+\b', line))
        for word in line_words:
            if word in difference_words:
                if word not in word_positions:
                    word_positions[word] = []
                word_positions[word].append({
                    'Slide': slide_number,
                    'Line': line_number,
                    'Position': line_words.index(word),
                    'LineContent': line
                })

    output = f'{input_pptx_file}_pptcomparison.txt'
    with open(output, 'w', encoding='utf-8') as result_file:
        for word, positions in word_positions.items():
            for data in positions:
                result_file.write(f"Word: {word}\n")
                result_file.write(f"Slide: {data['Slide']}\n")
                result_file.write(f"Line: {data['Line']}\n")
                result_file.write(f"Position: {data['Position']}\n")
                result_file.write(f"Line Content: {data['LineContent']}\n\n")

    logging.info(f"Words in PPT but not in HTML, along with positions, saved to {output}")

    # Comparison using BERT Cosine Similarity
    model_name = "bert-base-uncased"
    tokenizer = BertTokenizer.from_pretrained(model_name)
    model = BertModel.from_pretrained(model_name)

    tokens1 = tokenizer(ppt_text, return_tensors='pt', padding=True, truncation=True)
    tokens2 = tokenizer(html_text, return_tensors='pt', padding=True, truncation=True)

    with torch.no_grad():
        embeddings1 = model(**tokens1).last_hidden_state.mean(dim=1)
        embeddings2 = model(**tokens2).last_hidden_state.mean(dim=1)

    similarity = cosine_similarity(embeddings1, embeddings2)
    logging.info(f"Bert Cosine Similarity: {similarity[0][0]:.2f}")


# Example usage:
input_file_path = input('Enter the file path you want to process: ')
html_file_path = input('Enter the HTML file path for comparison: ')

if not os.path.exists(html_file_path):
    print(f'The HTML file {html_file_path} does not exist.')
else:
    html_output = 'output_html.txt'
    extract_text_from_html(html_file_path, html_output)
    with open(html_output, 'r', encoding='utf-8') as html_file:
        html_text = html_file.read()
    print(f"Extracted text content for HTML saved to {html_output}")

if not os.path.exists(input_file_path):
    print(f'The file {input_file_path} does not exist.')
else:
    file_extension = os.path.splitext(input_file_path)[1].lower()

    if file_extension == '.pdf':
        input_pdf_file = os.path.splitext(input_file_path)[0].strip()
        output_text_file = f"{input_pdf_file}_extracted.txt"
        extract_text_from_pdf(input_file_path, output_text_file)
        with open(output_text_file, 'r', encoding='utf-8') as pdf_file:
            pdf_text = pdf_file.read().lower()
        compare_pdf_with_html(pdf_text, html_text,input_pdf_file)

    elif file_extension == '.xlsx':
        input_excel_file = os.path.splitext(input_file_path)[0].strip()
        output_text_file = f"{input_excel_file}_extracted.txt"
        output_excel_file = extract_text_from_excel(input_file_path)
        ensure_directory_exists(output_excel_file)
        with open(output_excel_file, 'r', encoding='utf-8') as excel_file:
            excel_text = excel_file.read().lower()
        compare_excel_with_html(excel_text, html_text,input_excel_file)

    elif file_extension == '.docx':
        input_docs_file = os.path.splitext(input_file_path)[0].strip()
        output_text_file = f"{input_docs_file}_extracted.txt"
        output_docs_file = extract_text_from_docx(input_file_path)
        with open(output_docs_file, 'r', encoding='utf-8') as docx_file:
            docx_text = docx_file.read().lower()
        compare_docx_with_html(docx_text, html_text,input_docs_file)

    elif file_extension == '.pptx':
        input_pptx_file = os.path.splitext(input_file_path)[0].strip()
        output_text_file = f"{input_pptx_file}_extracted.txt"
        ensure_directory_exists('pptx_images')
        ensure_directory_exists('output_ppt')
        ppt_text = extract_text_from_ppt(input_file_path)
        compare_ppt_with_html(ppt_text, html_text,input_pptx_file)

    else:
        print(f"Unsupported file extension: {file_extension}")
import os
import logging
import nltk
import re
import torch
from transformers import BertTokenizer, BertModel
from sklearn.metrics.pairwise import cosine_similarity
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pytesseract
from PIL import Image
from inscriptis import get_text

def extract_text_from_ppt(pptx_file, input_html_file):
    # Process HTML
    if not os.path.exists(input_html_file):
        print(f'The file {input_html_file} does not exist.')
        exit()

    with open(input_html_file, 'r', encoding="ISO-8859-1") as file:
        html_content = file.read()
    text_content= get_text(html_content)

    def extract_table(table):
        # Function to extract text from a PowerPoint table
        data = []
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                cell_text = ""
                for paragraph in cell.text_frame.paragraphs:
                    cell_text += " ".join(run.text for run in paragraph.runs)
                row_text.append(cell_text)
            data.append("\n".join(row_text))
        return "\n".join(data)

    def extract_text_from_image(image_path):
        # Function to extract text from an image using Tesseract OCR
        try:
            image = Image.open(image_path)
            extracted_text = pytesseract.image_to_string(image)
            return extracted_text
        except Exception as e:
            print(f"Error while extracting text from image: {e}")
            return ""

    extracted_text = ""
    prs = Presentation(pptx_file)
    input_pptx_file = os.path.splitext(pptx_file)[0].strip()
    output_folder = f'{input_pptx_file}_images'

    for slide_number, slide in enumerate(prs.slides):
        extracted_text += f"Slide {slide_number + 1}:\n"

        for shape in slide.shapes:
            if shape.has_table:
                table_text = extract_table(shape.table)
                extracted_text += table_text + '\n\n'
            elif hasattr(shape, "text"):
                extracted_text += shape.text + '\n\n'
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_part = shape.image
                if image_part:
                    image_bytes = image_part.blob
                    image_filename = os.path.join(output_folder, f'image_slide_{slide_number + 1}.jpg')
                    with open(image_filename, 'wb') as f:
                        f.write(image_bytes)
                    extracted_image_text = extract_text_from_image(image_filename)
                    extracted_text += extracted_image_text + '\n\n'

    output_text_file = f'{input_pptx_file}_ppt_text.txt'
    with open(output_text_file, 'w', encoding='utf-8') as text_file:
        text_file.write(extracted_text)

        
    return extracted_text,text_content

def compare_ppt_with_html(ppt_text, html_text):
    try:
        ppt_words = set(nltk.word_tokenize(ppt_text))
        html_words = set(nltk.word_tokenize(html_text))

        # Finding the difference
        difference_words = ppt_words - html_words

        # Finding the line, position, page
        word_positions = {}
        ppt_lines = ppt_text.splitlines()

        page_number = 0
        line_number = 0

        for line in ppt_lines:
            if line.startswith("Slide"):
                match = re.match(r'Slide (\d+)', line)
                if match:
                    page_number = int(match.group(1))
                line_number = 0
            else:
                line_number += 1

            line_words = list(re.findall(r'\b\w+\b', line))
            for word in line_words:
                if word in difference_words:
                    if word not in word_positions:
                        word_positions[word] = []
                    word_positions[word].append({
                        'Page': page_number,
                        'Line': line_number,
                        'Position': line_words.index(word),
                        'LineContent': line
                    })

        output = f"ppt_htmlcompare.txt"
        with open(output, 'w', encoding='utf-8') as result_file:
            for word, positions in word_positions.items():
                for data in positions:
                    result_file.write(f"Word: {word}\n")
                    result_file.write(f"Page: {data['Page']}\n")
                    result_file.write(f"Line: {data['Line']}\n")
                    result_file.write(f"Position: {data['Position']}\n")
                    result_file.write(f"Line Content: {data['LineContent']}\n\n")

        logging.info(f"Words in PDF but not in HTML, along with positions, saved to {output}")

        # Comparison using Jaccard Similarity
        jaccard_similarity = len(ppt_words.intersection(html_words)) / len(ppt_words.union(html_words))
        percentage_difference = (1 - jaccard_similarity) * 100
        logging.info(f"Jaccard Similarity Score: {jaccard_similarity:.2f}")

        # Comparison using BERT Cosine Similarity
        model_name = "bert-base-uncased"
        tokenizer = BertTokenizer.from_pretrained(model_name)
        model = BertModel.from_pretrained(model_name)

        tokens1 = tokenizer(ppt_text, return_tensors='pt', padding=True, truncation=True)
        tokens2 = tokenizer(html_text, return_tensors='pt', padding=True, truncation=True)

        with torch.no_grad():
            embeddings1 = model(**tokens1).last_hidden_state.mean(dim=1)
            embeddings2 = model(**tokens2).last_hidden_state.mean(dim=1)

        similarity = cosine_similarity(embeddings1, embeddings2)
        bert_cosine_similarity = similarity[0][0].item()  # Convert float32 to a standard Python float
        logging.info(f"Bert Cosine Similarity: {bert_cosine_similarity:.2f}")

        # Include extracted PDF and HTML texts and comparison output in the response
        with open(output, 'r', encoding='utf-8') as result_file:
            output_content = result_file.read()

        response_data = {
            "bert_cosine_similarity": float(similarity[0][0]),
            "ppt_text": ppt_text,
            "html_text": html_text,
            "comparison_output": {
                "file_path": output,
                "content": output_content
            }
        }

        return response_data
    
    except Exception as e:
        return {"error": str(e)}